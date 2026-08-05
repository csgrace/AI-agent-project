"""PPT file parser using python-pptx.

Extracts structured content from PowerPoint files including titles,
body text, bullet points, speaker notes, tables, and images.
"""

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches


class PPTParser:
    """Parser for PowerPoint (.pptx/.ppt) files."""

    @staticmethod
    def parse(file_path: str | Path) -> dict[str, Any]:
        """Parse a .pptx file and return structured content.

        Args:
            file_path: Path to the .pptx file.

        Returns:
            JSON-compatible dict with file_name, total_slides, and slides list.
            Each slide contains slide_number, title, content, bullet_points,
            speaker_notes, tables, and images.

        Raises:
            FileNotFoundError: If the file does not exist.
            ValueError: If the file format is not supported.
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = path.suffix.lower()
        if ext not in (".pptx", ".ppt"):
            raise ValueError(f"Unsupported file format: {ext}. Only .pptx and .ppt are supported.")

        prs = Presentation(str(path))
        result: dict[str, Any] = {
            "file_name": path.name,
            "total_slides": len(prs.slides),
            "slides": [],
        }

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_data = _extract_slide(slide, slide_num)
            result["slides"].append(slide_data)

        return result


def _extract_slide(slide, slide_number: int) -> dict[str, Any]:
    """Extract data from a single slide."""
    title = ""
    content: list[str] = []
    bullet_points: list[str] = []
    tables: list[list[list[str]]] = []
    images: list[dict[str, str]] = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue

                if _is_title(shape, paragraph):
                    title = text
                elif paragraph.level > 0:
                    bullet_points.append(text)
                else:
                    content.append(text)

        if shape.has_table:
            table_data: list[list[str]] = []
            for row in shape.table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)

        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            image_info = _extract_image_info(shape, slide_number)
            if image_info:
                images.append(image_info)

    speaker_notes = _extract_speaker_notes(slide)

    # Deduplicate: remove bullet points that already appear in content
    content_set = set(content)
    bullet_points = [b for b in bullet_points if b not in content_set]

    return {
        "slide_number": slide_number,
        "title": title,
        "content": content,
        "bullet_points": bullet_points,
        "speaker_notes": speaker_notes,
        "tables": tables,
        "images": images,
    }


def _is_title(shape, paragraph) -> bool:
    """Heuristic: check if a paragraph is a slide title."""
    if shape.shape_id == 0:
        return True
    try:
        if paragraph.font.size and paragraph.font.size >= Inches(0.4):
            return True
    except Exception:
        pass
    return False


def _extract_speaker_notes(slide) -> str:
    """Extract speaker notes from a slide."""
    try:
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text.strip()
        return notes_text
    except Exception:
        return ""


def _extract_image_info(shape, slide_number: int) -> dict[str, str] | None:
    """Extract image path and caption from a shape."""
    try:
        image = shape.image
        ext = image.content_type.split("/")[-1] if "/" in image.content_type else "png"
        return {
            "path": f"slide_{slide_number}_img_{shape.shape_id}.{ext}",
            "caption": "",
        }
    except Exception:
        return None
